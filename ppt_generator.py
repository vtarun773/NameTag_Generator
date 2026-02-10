def generate_ppt(
    csv_path,
    template_path,
    image_dir,
    output_path,
    email_suffix
):

    import os
    import csv
    import tempfile
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from PIL import Image
    from pptx.enum.shapes import PP_PLACEHOLDER

    import csv

    people = []
    '''
    for filename in os.listdir(image_dir):
        name = os.path.join(image_dir, filename)
        new_name = name.lower()
        new_path = os.path.join(image_dir, new_name)
        
        if name != new_path:
            os.rename(name, new_name)
    
    '''
    encodings = ("utf-8-sig", "utf-8")
    last_error = None

    for enc in encodings:
        try:
            with open(csv_path, newline="", encoding=enc) as f:
                reader = csv.DictReader(f)
                for row in reader:
                    people.append({
                        "name": row["name"].strip(),
                        "email": row["ID"].strip(),
                        "department": row["department"].strip()
                    })
            break  # success → exit encoding loop

        except UnicodeDecodeError as e:
            people.clear()
            last_error = e

    else:
        raise RuntimeError(
            "Failed to read CSV file.\n"
            "Please save it as UTF-8 or UTF-8 with BOM (Excel default)."
        ) from last_error

    prs = Presentation(template_path)
    layout = prs.slide_layouts[0]

    FIELD_IDX = {}
    for ph in layout.placeholders:
        if ph.name == "NamePH":
            FIELD_IDX["name"] = ph.placeholder_format.idx
        elif ph.name == "EmailPH":
            FIELD_IDX["email"] = ph.placeholder_format.idx
        elif ph.name == "DeptPH":
            FIELD_IDX["department"] = ph.placeholder_format.idx
        elif ph.name == "ImagePH":
            FIELD_IDX["image"] = ph.placeholder_format.idx

    def set_text(ph, text, size, bold, color, align, font_name):
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color
        p.alignment = align

    def normalize_image(image_path):
        """
        Converts image into a PMO-safe RGB JPG.
        Returns new image path or None if unusable.
        """
        try:
            img = Image.open(image_path)

            img.verify()

            img = Image.open(image_path)
            img = img.convert("RGB")

            tmp = tempfile.NamedTemporaryFile(
                suffix=".jpg",
                delete=False
            )

            img.save(
                tmp.name,
                format="JPEG",
                quality=92,
                subsampling=0,
                optimize=True
            )
            return tmp.name

        except Exception as e:
            print(f"⚠ Bad image skipped: {os.path.basename(image_path)} ({e})")
            return None

    def add_image_to_placeholder(slide, placeholder_idx, image_dir, image_name_no_ext):
        """
        slide: pptx slide
        placeholder_idx: index of placeholder
        image_dir: folder where images are stored
        image_name_no_ext: ID only (without extension)
        """
        ph = slide.placeholders[placeholder_idx]

        IMG_EXTS = [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff"]
        raw_image = None

        # find first available image with supported extension
        for ext in IMG_EXTS:
            temp_path = os.path.join(image_dir, image_name_no_ext + ext)
            if os.path.exists(temp_path):
                raw_image = temp_path
                break

        if not raw_image:
            print(f"⚠ Image not found for ID '{image_name_no_ext}' in {image_dir}")
            return

        safe_image = normalize_image(raw_image)
        if not safe_image:
            return  # placeholder stays empty

        # Placeholder geometry
        left = ph.left
        top = ph.top
        width = ph.width
        height = ph.height

        # Remove placeholder
        ph._element.getparent().remove(ph._element)

        # Add picture stretched to placeholder
        try:
            slide.shapes.add_picture(safe_image, left, top, width=width, height=height)
        except Exception as e:
            print(f"⚠ Image skipped for '{image_name_no_ext}': {e}")

    for person in people:
        slide = prs.slides.add_slide(layout)

        set_text(
            slide.placeholders[FIELD_IDX["name"]],
            person["name"],
            32,
            True,
            RGBColor(255, 255, 255),
            PP_ALIGN.CENTER,
            "Montserrat ExtraBold"
        )

        set_text(
            slide.placeholders[FIELD_IDX["email"]],
            person["email"] + email_suffix,
            20,
            False,
            RGBColor(5, 77, 162),
            PP_ALIGN.LEFT,
            "Montserrat ExtraBold"
        )

        set_text(
            slide.placeholders[FIELD_IDX["department"]],
            person["department"],
            20,
            False,
            RGBColor(5, 77, 162),
            PP_ALIGN.LEFT,
            "Montserrat ExtraBold"
        )

        add_image_to_placeholder(
            slide,
            FIELD_IDX["image"],
            image_dir,
            person["email"]  # just the ID, no extension
        )

    prs.save(output_path)



